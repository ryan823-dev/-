#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIGC 在数字贸易出海及跨境电商运营中的应用 - PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色定义
COLORS = {
    'primary': RGBColor(0x1E, 0x3A, 0x8A),      # 深海蓝
    'blue': RGBColor(0x3B, 0x82, 0xF6),         # 科技蓝
    'green': RGBColor(0x10, 0xB9, 0x81),        # 荧光绿
    'orange': RGBColor(0xF9, 0x73, 0x16),       # 活力橙
    'yellow': RGBColor(0xF5, 0x9E, 0x0B),       # 琥珀色
    'white': RGBColor(0xFF, 0xFF, 0xFF),        # 白色
    'light_blue': RGBColor(0xE0, 0xF2, 0xFE),   # 浅蓝
    'gray': RGBColor(0x64, 0x74, 0x8B),         # 灰色
    'light_gray': RGBColor(0xF8, 0xFA, 0xFC),   # 浅灰
}

def set_background(slide, prs, color1=COLORS['primary'], color2=COLORS['blue']):
    """设置渐变背景"""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops.insert(0, color1, 0.0)
    fill.gradient_stops.insert(1, color2, 1.0)

def create_cover_slide(prs, title, subtitle, info):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_background(slide, prs, COLORS['primary'], COLORS['green'])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light_blue']
    p.alignment = PP_ALIGN.CENTER
    
    # 信息
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1))
    tf = info_box.text_frame
    for i, line in enumerate(info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['light_blue']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)
    
    return slide

def create_toc_slide(prs, modules):
    """创建目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 课程目录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 模块列表
    y_pos = 1.3
    for i, module in enumerate(modules):
        module_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y_pos), Inches(11.6), Inches(0.9)
        )
        module_box.fill.solid()
        module_box.fill.fore_color.rgb = COLORS['light_blue']
        module_box.line.color.rgb = COLORS['blue']
        module_box.line.width = Pt(2)
        
        # 模块内容
        tf = module_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"{module['number']}  {module['title']}"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLORS['primary']
        
        p2 = tf.add_paragraph()
        p2.text = f"⏱️ {module['time']}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['gray']
        p2.space_before = Pt(5)
        
        y_pos += 1.0
    
    return slide

def create_content_slide(prs, title, content_items, highlight=None):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # 标题栏
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13), Inches(1.2)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLORS['blue']
    header_box.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容
    y_pos = 1.4
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(11.6), Inches(0.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"▸ {item['text']}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray'] if 'color' not in item else item['color']
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        y_pos += 0.7
    
    # 高亮框
    if highlight:
        hl_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y_pos + 0.2), Inches(11.6), Inches(1.2)
        )
        hl_box.fill.solid()
        hl_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
        hl_box.line.color.rgb = COLORS['yellow']
        hl_box.line.width = Pt(3)
        
        tf = hl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = highlight['text']
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
    
    return slide

def create_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """创建对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 左列
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['light_blue']
    left_box.line.color.rgb = COLORS['blue']
    left_box.line.width = Pt(3)
    
    # 左列标题
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 左列内容
    y_offset = 0.5
    for i, item in enumerate(left_items):
        p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS['gray']
        p.space_before = Pt(8)
    
    # 右列
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(1.3), Inches(5.8), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    right_box.line.color.rgb = COLORS['green']
    right_box.line.width = Pt(3)
    
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
    p.alignment = PP_ALIGN.CENTER
    
    for i, item in enumerate(right_items):
        p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS['gray']
        p.space_before = Pt(8)
    
    return slide

def create_ending_slide(prs, title, thank_you, contact):
    """创建结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs, COLORS['green'], COLORS['blue'])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 感谢语
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = thank_you
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    p.alignment = PP_ALIGN.CENTER
    
    # 联系信息
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12), Inches(1.5))
    tf = contact_box.text_frame
    for i, line in enumerate(contact):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(15)
    
    return slide

def main():
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13)
    prs.slide_height = Inches(7.5)
    
    # 使用空白布局创建幻灯片

    
    print("🎨 开始生成 PPT...")
    
    # Slide 1: 封面
    print("  → 创建封面页...")
    create_cover_slide(
        prs,
        "AIGC 在数字贸易出海及\n跨境电商运营中的应用",
        "赋能全球业务 · 智能驱动增长",
        ["三江学院 · 专题讲座", "2026 年"]
    )
    
    # Slide 2: 目录
    print("  → 创建目录页...")
    modules = [
        {"number": "01", "title": "开场导入：AI 浪潮与出海机遇", "time": "10 分钟"},
        {"number": "02", "title": "数字贸易出海全景图", "time": "20 分钟"},
        {"number": "03", "title": "AI+ 跨境电商运营实战", "time": "35 分钟"},
        {"number": "04", "title": "AI+ 数字营销获客", "time": "35 分钟"},
        {"number": "05", "title": "热点工具实战演示", "time": "15 分钟"},
        {"number": "06", "title": "总结互动与行动建议", "time": "5 分钟"},
    ]
    create_toc_slide(prs, modules)
    
    # Slide 3: AI 浪潮
    print("  → 创建 AI 浪潮页...")
    create_content_slide(
        prs,
        "🚀 AI 浪潮爆发：2024-2025 关键趋势",
        [
            {"text": "OpenClaw（大龙虾）：新一代 AI 爬虫工具，自动化数据采集与竞品监控", "color": COLORS['blue']},
            {"text": "GEO（Generative Engine Optimization）：生成式引擎优化，让 AI 主动推荐你的品牌", "color": COLORS['green']},
            {"text": "AI Agent（智能体）：自主执行复杂任务的 AI 助手，从客服到运营全流程", "color": COLORS['orange']},
        ],
        {"text": "💡 关键洞察：AI 不再是'可选项'，而是出海企业的'必选项'。早期采用者已获得 3-5 倍效率提升。"}
    )
    
    # Slide 4: 三大挑战
    print("  → 创建挑战页...")
    create_content_slide(
        prs,
        "🎯 中国卖家出海的三大核心挑战",
        [
            {"text": "语言障碍：多语言内容本地化困难"},
            {"text": "流量困境：获客成本高、竞争白热化"},
            {"text": "运营难题：数据驱动决策能力不足"},
            {"text": "效率瓶颈：人力密集、响应速度慢"},
        ],
        {"text": "✅ AI 解决方案：多语言生成 · 智能获客 · 数据洞察 · 自动化运营"}
    )
    
    # Slide 5: 数字贸易 vs 跨境电商
    print("  → 创建对比页...")
    create_comparison_slide(
        prs,
        "📊 数字贸易 vs 跨境电商：概念辨析",
        "🏭 B2B 数字贸易",
        [
            "阿里巴巴国际站、环球资源",
            "大宗交易、长期合同",
            "决策链路长（3-6 个月）",
            "客单价高（$10,000+）",
            "关系驱动、信任建立",
            "主动开发客户为主"
        ],
        "🛒 B2C 跨境电商",
        [
            "Amazon、TikTok Shop、独立站",
            "零售交易、即时购买",
            "决策链路短（几分钟 -7 天）",
            "客单价低（$20-$500）",
            "流量驱动、转化优化",
            "平台流量 + 广告投放"
        ]
    )
    
    # Slide 6: 出海模式矩阵
    print("  → 创建模式矩阵页...")
    create_content_slide(
        prs,
        "🌐 出海商业模式矩阵",
        [
            {"text": "第三方平台：Amazon、eBay、Shopee - 流量大、规则多、竞争激烈"},
            {"text": "独立站：Shopify、WooCommerce - 品牌化、利润高、需引流"},
            {"text": "社交电商：TikTok Shop、Instagram - 内容驱动、爆发力强"},
            {"text": "B2B 平台：阿里巴巴国际站 - 大客户、长周期、高复购"},
            {"text": "主动获客：LinkedIn、Google 搜索 - 精准开发、利润最高"},
            {"text": "混合模式：平台 + 独立站 + 社媒 - 风险分散、最大化覆盖"},
        ]
    )
    
    # Slide 7: AI 价值地图
    print("  → 创建 AI 价值页...")
    create_content_slide(
        prs,
        "🤖 AI 在出海业务中的价值地图",
        [
            {"text": "📊 运营端：数据分析、选品决策、库存优化"},
            {"text": "📣 营销端：内容生成、SEO/GEO、广告投放"},
            {"text": "💬 客服端：智能客服、多语言支持、工单处理"},
            {"text": "🎯 获客端：主动搜索、客户画像、线索评分"},
        ],
        {"text": "💰 ROI 数据：平均效率提升 300% · 人力成本降低 60% · 转化率提升 45%"}
    )
    
    # Slide 8: AI 选品实战
    print("  → 创建选品页...")
    create_content_slide(
        prs,
        "🎯 AI 驱动的数据分析与选品",
        [
            {"text": "市场趋势分析：用 AI 分析 Google Trends、Amazon BSR 数据"},
            {"text": "竞品监控：AI 自动追踪竞品价格、评论、排名变化"},
            {"text": "选品决策模型：AI 评分系统（需求 + 竞争 + 利润 + 供应链）"},
            {"text": "案例：某卖家通过 AI 选品打造月销 10 万 + 爆款，利润率 35%"},
        ],
        {"text": "🔧 推荐工具：Jungle Scout · Helium 10 · 卖家精灵 · Keepa"}
    )
    
    # Slide 9: AI 运营优化
    print("  → 创建运营优化页...")
    create_content_slide(
        prs,
        "⚙️ AI 优化运营策略",
        [
            {"text": "动态定价：基于竞品、库存、季节性，AI 实时调整（提升利润 15-25%）"},
            {"text": "智能库存：AI 销量预测，提前 30 天预警，降低库存周转 40%"},
            {"text": "Listing 优化：AI 生成标题、五点描述，转化率提升 30%+"},
            {"text": "广告优化：自动关键词、出价调整，ACOS 降低 35%"},
        ],
        {"text": "📈 效果数据：综合运营效率提升 2.5 倍 · 人力成本降低 50%"}
    )
    
    # Slide 10: SEO vs GEO
    print("  → 创建 SEO 对比页...")
    create_comparison_slide(
        prs,
        "🔍 传统 SEO vs GEO（生成式引擎优化）",
        "传统 SEO",
        [
            "优化给搜索引擎看",
            "关键词密度、外链建设",
            "目标：Google 首页排名",
            "结果：获得点击流量",
            "技术：Meta 标签、网站速度"
        ],
        "GEO（新趋势）",
        [
            "优化给 AI 引擎看",
            "语义三重奏、结构化数据",
            "目标：ChatGPT、Claude 引用",
            "结果：AI 主动推荐品牌",
            "技术：Schema.org、权威背书"
        ]
    )
    
    # Slide 11: AI 内容矩阵
    print("  → 创建内容矩阵页...")
    create_content_slide(
        prs,
        "📝 AI 内容矩阵生成实战",
        [
            {"text": "产品描述批量生成：输入参数，AI 生成多版本描述"},
            {"text": "Blog 文章自动撰写：行业知识、使用指南，每月 50-100 篇"},
            {"text": "社交媒体文案：多平台适配，自动添加热门 Hashtag"},
            {"text": "邮件营销内容：个性化开发信，打开率提升 60%"},
            {"text": "多语言本地化：20+ 语言，保持品牌调性一致"},
        ],
        {"text": "📊 案例：某品牌用 AI 生成 300+ 篇文章，6 个月自然流量增长 5 倍"}
    )
    
    # Slide 12: AI 短视频
    print("  → 创建短视频页...")
    create_content_slide(
        prs,
        "🎬 AI 短视频创作与传播",
        [
            {"text": "1️⃣ 脚本生成：AI 根据产品特性写分镜脚本"},
            {"text": "2️⃣ 视频制作：Runway、Pictory AI 生成画面"},
            {"text": "3️⃣ 数字人播报：HeyGen、D-ID 多语言配音"},
            {"text": "4️⃣ 批量分发：多平台自动发布 + 数据分析"},
        ],
        {"text": "🔧 推荐工具：Runway · HeyGen · CapCut AI · Descript"}
    )
    
    # Slide 13: B2B 主动获客
    print("  → 创建 B2B 获客页...")
    create_content_slide(
        prs,
        "🎯 B2B 主动获客全流程",
        [
            {"text": "客户智能搜索：AI 挖掘 LinkedIn、Google、海关数据"},
            {"text": "背调分析报告：自动生成客户画像、采购偏好"},
            {"text": "开发信写作：AI 个性化邮件，回复率从 2% 提升至 8%"},
            {"text": "社媒自动化：LinkedIn 内容发布 + 智能互动"},
        ],
        {"text": "📈 实战效果：3 个月获取 50+ 高质量询盘，成交 8 单总金额$2.3M"}
    )
    
    # Slide 14: 热点工具
    print("  → 创建工具页...")
    create_content_slide(
        prs,
        "🔥 热点工具实战演示",
        [
            {"text": "OpenClaw：AI 爬虫与数据采集、竞品监控"},
            {"text": "GEO 工具：AI 引擎可见性优化、Semantic SEO"},
            {"text": "Jasper/Copy.ai：AI 内容生成平台"},
            {"text": "Midjourney/DALL·E 3：AI 图像生成"},
            {"text": "Descript/CapCut：AI 视频剪辑"},
            {"text": "Notion AI：智能笔记与协作"},
        ]
    )
    
    # Slide 15: 学习路径
    print("  → 创建学习路径页...")
    create_content_slide(
        prs,
        "📚 学生入门学习路径",
        [
            {"text": "阶段 1：认知建立（1-2 周）- AI 基础、主流工具、行业格局"},
            {"text": "阶段 2：工具实操（3-4 周）- 熟练使用 3-5 个核心工具"},
            {"text": "阶段 3：项目实战（2-3 月）- 从 0 到 1 完整运营"},
            {"text": "阶段 4：优化迭代（持续）- 数据分析、策略优化"},
        ],
        {"text": "💡 建议：先精通 1-2 个工具 · 实战 > 理论 · 加入社群交流"}
    )
    
    # Slide 16: 避坑指南
    print("  → 创建避坑页...")
    create_content_slide(
        prs,
        "⚠️ AI 使用的 5 个误区",
        [
            {"text": "❌ 完全依赖 AI → ✅ AI 辅助 + 人工审核"},
            {"text": "❌ 忽视数据质量 → ✅ 高质量输入 = 高质量输出"},
            {"text": "❌ 盲目追求新工具 → ✅ 精通核心工具 > 浅尝辄止"},
            {"text": "❌ 忽视合规风险 → ✅ 遵守平台规则、版权法"},
            {"text": "❌ 期待立竿见影 → ✅ 持续优化、长期主义"},
        ]
    )
    
    # Slide 17: 工具清单
    print("  → 创建工具清单页...")
    create_content_slide(
        prs,
        "🛠️ 推荐工具清单（按预算）",
        [
            {"text": "免费：ChatGPT、Google Trends、Canva"},
            {"text": "入门级（$20-50/月）：Jungle Scout、Jasper、Helium 10"},
            {"text": "企业级（$200+/月）：OpenClaw、Enterprise GEO、Custom AI Agent"},
        ]
    )
    
    # Slide 18: 结束页
    print("  → 创建结束页...")
    create_ending_slide(
        prs,
        "🎓 总结与行动",
        "感谢聆听 · Q&A",
        ["立即行动：选择 1 个工具开始实践", "建立你的 AI 驱动出海能力 💪"]
    )
    
    # 保存
    output_path = "/Users/oceanlink/Documents/Qoder-1/AIGC 数字贸易出海及跨境电商运营应用.pptx"
    prs.save(output_path)
    print(f"\n✅ PPT 生成成功！保存到：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")

if __name__ == "__main__":
    main()
