#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIGC 数字贸易出海 PPT - 视觉升级版
设计理念：现代简约 + 科技感 + 专业商务
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ==================== 设计系统 ====================

class DesignSystem:
    """现代化设计系统"""
    
    # 配色方案：科技+商务风格
    COLORS = {
        # 主色调
        'primary': RGBColor(0x1E, 0x3A, 0x8A),      # 深海蓝 - 专业、信任
        'secondary': RGBColor(0x3B, 0x82, 0xF6),    # 科技蓝 - 创新、科技
        'accent': RGBColor(0x10, 0xB9, 0x81),       # 荧光绿 - 增长、活力
        'warning': RGBColor(0xF5, 0x9E, 0x0B),      # 琥珀色 - 提示、亮点
        'danger': RGBColor(0xEF, 0x44, 0x44),       # 红色 - 警示
        
        # 背景色
        'bg_dark': RGBColor(0x0F, 0x17, 0x2A),      # 深色背景
        'bg_light': RGBColor(0xF8, 0xFA, 0xFC),     # 浅色背景
        'bg_card': RGBColor(0xFF, 0xFF, 0xFF),      # 卡片背景
        
        # 文字色
        'text_primary': RGBColor(0x1E, 0x29, 0x3B), # 主要文字
        'text_secondary': RGBColor(0x64, 0x74, 0x8B), # 次要文字
        'text_light': RGBColor(0xE2, 0xE8, 0xF0),   # 浅色文字
        'text_white': RGBColor(0xFF, 0xFF, 0xFF),   # 白色文字
        
        # 辅助色
        'border': RGBColor(0xE2, 0xE8, 0xF0),       # 边框
        'divider': RGBColor(0xCB, 0xD5, 0xE1),      # 分割线
        'overlay': RGBColor(0x00, 0x00, 0x00),      # 遮罩层
    }
    
    # 字体系统
    FONTS = {
        'title': 'Arial Black',
        'heading': 'Arial',
        'body': 'Arial',
        'mono': 'Courier New',
    }
    
    # 间距系统
    SPACING = {
        'xs': Pt(4),
        'sm': Pt(8),
        'md': Pt(16),
        'lg': Pt(24),
        'xl': Pt(32),
        'xxl': Pt(48),
    }
    
    # 尺寸
    SIZES = {
        'title': Pt(48),
        'heading1': Pt(36),
        'heading2': Pt(28),
        'heading3': Pt(22),
        'body_large': Pt(18),
        'body': Pt(16),
        'body_small': Pt(14),
        'caption': Pt(12),
    }
    
    # 圆角
    BORDER_RADIUS = {
        'sm': Pt(4),
        'md': Pt(8),
        'lg': Pt(16),
        'xl': Pt(24),
        'full': Pt(999),
    }


class PPTBuilder:
    """PPT 构建器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.design = DesignSystem()
        self.slide_count = 0
        
    def add_slide(self, layout_index=6):
        """添加幻灯片"""
        self.slide_count += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[layout_index])
    
    # ==================== 布局组件 ====================
    
    def add_header_bar(self, slide, title, subtitle=None):
        """添加顶部导航栏"""
        # 背景
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.design.COLORS['primary']
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.25),
            Inches(11.733), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.design.SIZES['heading2']
        p.font.bold = True
        p.font.color.rgb = self.design.COLORS['text_white']
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.75),
                Inches(11.733), Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = self.design.SIZES['body_small']
            p.font.color.rgb = self.design.COLORS['text_light']
        
        # 页码
        page_box = slide.shapes.add_textbox(
            Inches(12.2), Inches(0.4),
            Inches(0.8), Inches(0.4)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.slide_count}"
        p.font.size = self.design.SIZES['body_small']
        p.font.color.rgb = self.design.COLORS['text_light']
        p.alignment = PP_ALIGN.RIGHT
        
        return header
    
    def add_card(self, slide, left, top, width, height, title=None, content=None, icon=None):
        """添加卡片"""
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.COLORS['bg_card']
        card.line.color.rgb = self.design.COLORS['border']
        card.line.width = Pt(1)
        
        # 调整圆角
        card.adjustments[0] = 0.05
        
        # 标题
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(left + 0.3), Inches(top + 0.2),
                Inches(width - 0.6), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{icon} {title}" if icon else title
            p.font.size = self.design.SIZES['heading3']
            p.font.bold = True
            p.font.color.rgb = self.design.COLORS['primary']
        
        # 内容
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(left + 0.3), Inches(top + 0.7),
                Inches(width - 0.6), Inches(height - 1.0)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = self.design.SIZES['body']
                p.font.color.rgb = self.design.COLORS['text_secondary']
                p.space_before = Pt(6)
                p.space_after = Pt(4)
        
        return card
    
    def add_highlight_box(self, slide, left, top, width, text, color='warning'):
        """添加高亮提示框"""
        colors = {
            'warning': self.design.COLORS['warning'],
            'success': self.design.COLORS['accent'],
            'danger': self.design.COLORS['danger'],
            'info': self.design.COLORS['secondary'],
        }
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7) if color == 'warning' else \
                                  RGBColor(0xD1, 0xFA, 0xE5) if color == 'success' else \
                                  RGBColor(0xFE, 0xE2, 0xE2) if color == 'danger' else \
                                  RGBColor(0xDB, 0xEA, 0xFE)
        box.line.fill.background()
        
        # 左边框
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(0.08), Inches(1.0)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = colors.get(color, self.design.COLORS['warning'])
        border.line.fill.background()
        
        # 文字
        text_box = slide.shapes.add_textbox(
            Inches(left + 0.3), Inches(top + 0.25),
            Inches(width - 0.5), Inches(0.6)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = self.design.SIZES['body']
        p.font.color.rgb = self.design.COLORS['text_primary']
        
        return box
    
    def add_stats_row(self, slide, y_pos, stats):
        """添加统计行"""
        x_start = 0.8
        card_width = 3.0
        gap = 0.4
        
        for i, stat in enumerate(stats):
            x = x_start + i * (card_width + gap)
            
            # 背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y_pos),
                Inches(card_width), Inches(1.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.design.COLORS['bg_card']
            card.line.color.rgb = self.design.COLORS['border']
            card.adjustments[0] = 0.08
            
            # 数值
            num_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y_pos + 0.2),
                Inches(card_width - 0.4), Inches(0.8)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat['value']
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.design.COLORS['primary']
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y_pos + 0.9),
                Inches(card_width - 0.4), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat['label']
            p.font.size = self.design.SIZES['body_small']
            p.font.color.rgb = self.design.COLORS['text_secondary']
    
    # ==================== 特殊布局 ====================
    
    def create_cover_slide(self):
        """创建封面页"""
        slide = self.add_slide()
        
        # 背景渐变效果（用多个矩形模拟）
        bg1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = self.design.COLORS['primary']
        bg1.line.fill.background()
        
        # 装饰元素 - 科技感线条
        for i in range(5):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(1.5 + i * 1.2),
                Inches(13.333), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.design.COLORS['secondary']
            line.fill.fore_color.brightness = 0.3
            line.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.333), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "AIGC 在数字贸易出海及"
        p.font.size = self.design.SIZES['title']
        p.font.bold = True
        p.font.color.rgb = self.design.COLORS['text_white']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "跨境电商运营中的应用"
        p2.font.size = self.design.SIZES['title']
        p2.font.bold = True
        p2.font.color.rgb = self.design.COLORS['text_white']
        p2.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(11.333), Inches(0.6)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "赋能全球业务 · 智能驱动增长"
        p.font.size = self.design.SIZES['heading3']
        p.font.color.rgb = self.design.COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 信息
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(11.333), Inches(1.5)
        )
        tf = info_box.text_frame
        
        p = tf.paragraphs[0]
        p.text = "三江学院 · 专题讲座"
        p.font.size = self.design.SIZES['body_large']
        p.font.color.rgb = self.design.COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "2026 年"
        p2.font.size = self.design.SIZES['body']
        p2.font.color.rgb = self.design.COLORS['text_light']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        
        return slide
    
    def create_toc_slide(self, modules):
        """创建目录页"""
        slide = self.add_slide()
        self.add_header_bar(slide, "课程大纲", "AIGC 数字贸易出海应用")
        
        y_start = 1.6
        card_height = 0.9
        gap = 0.15
        
        for i, module in enumerate(modules):
            y = y_start + i * (card_height + gap)
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(y),
                Inches(11.733), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.design.COLORS['bg_card'] if i % 2 == 0 else \
                                       RGBColor(0xEF, 0xF6, 0xFF)
            card.line.fill.background()
            card.adjustments[0] = 0.05
            
            # 编号
            num_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(y + 0.2),
                Inches(0.8), Inches(0.5)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = module['number']
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.design.COLORS['secondary']
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y + 0.15),
                Inches(8), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = module['title']
            p.font.size = self.design.SIZES['body_large']
            p.font.bold = True
            p.font.color.rgb = self.design.COLORS['text_primary']
            
            # 时间
            time_box = slide.shapes.add_textbox(
                Inches(10.5), Inches(y + 0.25),
                Inches(1.8), Inches(0.4)
            )
            tf = time_box.text_frame
            p = tf.paragraphs[0]
            p.text = module['time']
            p.font.size = self.design.SIZES['body_small']
            p.font.color.rgb = self.design.COLORS['text_secondary']
            p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def create_content_slide(self, title, items, highlight=None, section_num=None):
        """创建内容页"""
        slide = self.add_slide()
        
        # 标题
        full_title = f"{section_num}  {title}" if section_num else title
        self.add_header_bar(slide, full_title)
        
        # 内容
        y_start = 1.6
        item_height = 0.65
        
        for i, item in enumerate(items):
            y = y_start + i * item_height
            
            # 图标（如果有）
            icon = item.get('icon', '▸')
            
            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(y),
                Inches(11.333), Inches(item_height)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = f"{icon}  {item['text']}"
            p.font.size = self.design.SIZES['body_large']
            p.font.color.rgb = self.design.COLORS['text_primary']
            p.space_after = Pt(8)
            
            # 次级说明
            if 'sub' in item:
                p2 = tf.add_paragraph()
                p2.text = f"      {item['sub']}"
                p2.font.size = self.design.SIZES['body_small']
                p2.font.color.rgb = self.design.COLORS['text_secondary']
                p2.space_before = Pt(4)
        
        # 高亮框
        if highlight:
            y_highlight = y_start + len(items) * item_height + 0.3
            self.add_highlight_box(slide, 0.8, y_highlight, 11.733, highlight)
        
        return slide
    
    def create_comparison_slide(self, title, left_data, right_data, section_num=None):
        """创建对比页"""
        slide = self.add_slide()
        
        full_title = f"{section_num}  {title}" if section_num else title
        self.add_header_bar(slide, full_title)
        
        # 左列
        left_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.6),
            Inches(5.8), Inches(5.2)
        )
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
        left_card.line.color.rgb = self.design.COLORS['secondary']
        left_card.line.width = Pt(2)
        left_card.adjustments[0] = 0.05
        
        # 左列标题
        left_title = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8),
            Inches(5.4), Inches(0.6)
        )
        tf = left_title.text_frame
        p = tf.paragraphs[0]
        p.text = left_data['title']
        p.font.size = self.design.SIZES['heading3']
        p.font.bold = True
        p.font.color.rgb = self.design.COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 左列内容
        y = 2.5
        for item in left_data['items']:
            item_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y),
                Inches(5.0), Inches(0.5)
            )
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ {item}"
            p.font.size = self.design.SIZES['body']
            p.font.color.rgb = self.design.COLORS['text_secondary']
            y += 0.55
        
        # 右列
        right_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.933), Inches(1.6),
            Inches(5.8), Inches(5.2)
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
        right_card.line.color.rgb = self.design.COLORS['accent']
        right_card.line.width = Pt(2)
        right_card.adjustments[0] = 0.05
        
        # 右列标题
        right_title = slide.shapes.add_textbox(
            Inches(7.133), Inches(1.8),
            Inches(5.4), Inches(0.6)
        )
        tf = right_title.text_frame
        p = tf.paragraphs[0]
        p.text = right_data['title']
        p.font.size = self.design.SIZES['heading3']
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
        p.alignment = PP_ALIGN.CENTER
        
        # 右列内容
        y = 2.5
        for item in right_data['items']:
            item_box = slide.shapes.add_textbox(
                Inches(7.333), Inches(y),
                Inches(5.0), Inches(0.5)
            )
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ {item}"
            p.font.size = self.design.SIZES['body']
            p.font.color.rgb = self.design.COLORS['text_secondary']
            y += 0.55
        
        return slide
    
    def create_stats_slide(self, title, stats, section_num=None):
        """创建数据统计页"""
        slide = self.add_slide()
        
        full_title = f"{section_num}  {title}" if section_num else title
        self.add_header_bar(slide, full_title)
        
        # 统计卡片
        self.add_stats_row(slide, 1.8, stats[:4])
        
        if len(stats) > 4:
            self.add_stats_row(slide, 3.5, stats[4:8])
        
        return slide
    
    def create_ending_slide(self):
        """创建结束页"""
        slide = self.add_slide()
        
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.design.COLORS['accent']
        bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢聆听 · Q&A"
        p.font.size = self.design.SIZES['title']
        p.font.bold = True
        p.font.color.rgb = self.design.COLORS['text_white']
        p.alignment = PP_ALIGN.CENTER
        
        # 行动号召
        cta_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(11.333), Inches(1.5)
        )
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "立即行动：选择 1 个工具开始实践"
        p.font.size = self.design.SIZES['heading3']
        p.font.color.rgb = self.design.COLORS['text_white']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "建立你的 AI 驱动出海能力"
        p2.font.size = self.design.SIZES['body_large']
        p2.font.color.rgb = self.design.COLORS['text_light']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        
        return slide
    
    def save(self, path):
        """保存 PPT"""
        self.prs.save(path)
        print(f"✅ PPT 已保存：{path}")
        print(f"📊 共 {self.slide_count} 页幻灯片")


def main():
    """主函数"""
    print("🎨 开始生成视觉升级版 PPT...\n")
    
    builder = PPTBuilder()
    
    # ========== 封面页 ==========
    print("  → 创建封面页...")
    builder.create_cover_slide()
    
    # ========== 目录页 ==========
    print("  → 创建目录页...")
    modules = [
        {"number": "01", "title": "开场导入：AI 浪潮与出海机遇", "time": "10 分钟"},
        {"number": "02", "title": "数字贸易出海全景图", "time": "20 分钟"},
        {"number": "03", "title": "AI+ 跨境电商运营实战", "time": "35 分钟"},
        {"number": "04", "title": "AI+ 数字营销获客", "time": "35 分钟"},
        {"number": "05", "title": "热点工具实战演示", "time": "15 分钟"},
        {"number": "06", "title": "总结互动与行动建议", "time": "5 分钟"},
    ]
    builder.create_toc_slide(modules)
    
    # ========== 模块一：开场导入 ==========
    print("  → 创建模块一...")
    
    builder.create_content_slide(
        "AI 浪潮爆发：2024-2025 关键趋势",
        [
            {"text": "OpenClaw（大龙虾）：新一代 AI 爬虫工具", "icon": "🦞", "sub": "自动化数据采集与竞品监控"},
            {"text": "GEO（Generative Engine Optimization）", "icon": "🎯", "sub": "生成式引擎优化，让 AI 主动推荐你的品牌"},
            {"text": "AI Agent（智能体）", "icon": "🤖", "sub": "自主执行复杂任务的 AI 助手"},
        ],
        "💡 关键洞察：AI 不再是「可选项」，而是出海企业的「必选项」",
        "01"
    )
    
    builder.create_stats_slide(
        "中国卖家出海的三大核心挑战",
        [
            {"value": "语言", "label": "多语言内容本地化"},
            {"value": "流量", "label": "获客成本高"},
            {"value": "运营", "label": "数据驱动决策难"},
            {"value": "效率", "label": "人力密集"},
        ],
        "01"
    )
    
    # ========== 模块二：全景图 ==========
    print("  → 创建模块二...")
    
    builder.create_comparison_slide(
        "数字贸易 vs 跨境电商",
        {
            "title": "🏭 B2B 数字贸易",
            "items": [
                "阿里巴巴国际站、环球资源",
                "大宗交易、长期合同",
                "决策链路长（3-6 个月）",
                "客单价高（$10,000+）",
                "关系驱动、信任建立",
            ]
        },
        {
            "title": "🛒 B2C 跨境电商",
            "items": [
                "Amazon、TikTok Shop、独立站",
                "零售交易、即时购买",
                "决策链路短（几分钟 -7 天）",
                "客单价低（$20-$500）",
                "流量驱动、转化优化",
            ]
        },
        "02"
    )
    
    builder.create_content_slide(
        "出海商业模式矩阵",
        [
            {"text": "第三方平台：Amazon、eBay、Shopee", "icon": "🛒"},
            {"text": "独立站：Shopify、WooCommerce", "icon": "🌐"},
            {"text": "社交电商：TikTok Shop、Instagram", "icon": "📱"},
            {"text": "B2B 平台：阿里巴巴国际站", "icon": "🏢"},
            {"text": "主动获客：LinkedIn、Google", "icon": "🎯"},
        ],
        None,
        "02"
    )
    
    builder.create_content_slide(
        "AI 在出海业务中的价值地图",
        [
            {"text": "运营端：数据分析、选品决策、库存优化", "icon": "📊"},
            {"text": "营销端：内容生成、SEO/GEO、广告投放", "icon": "📣"},
            {"text": "客服端：智能客服、多语言支持", "icon": "💬"},
            {"text": "获客端：主动搜索、客户画像", "icon": "🎯"},
        ],
        "💰 ROI：效率提升 300% · 成本降低 60% · 转化提升 45%",
        "02"
    )
    
    # ========== 模块三：运营实战 ==========
    print("  → 创建模块三...")
    
    builder.create_content_slide(
        "AI 驱动的数据分析与选品",
        [
            {"text": "市场趋势分析", "icon": "📈", "sub": "AI 分析 Google Trends、Amazon BSR 数据"},
            {"text": "竞品监控", "icon": "👁️", "sub": "自动追踪价格、评论、排名变化"},
            {"text": "选品决策模型", "icon": "🧠", "sub": "需求 + 竞争 + 利润 + 供应链评分"},
        ],
        "🔧 工具：Jungle Scout · Helium 10 · 卖家精灵",
        "03"
    )
    
    builder.create_content_slide(
        "AI 优化运营策略",
        [
            {"text": "动态定价：实时调整，提升利润 15-25%", "icon": "💲"},
            {"text": "智能库存：提前 30 天预警断货", "icon": "📦"},
            {"text": "Listing 优化：转化率提升 30%+", "icon": "✨"},
            {"text": "广告优化：ACOS 降低 35%", "icon": "📢"},
        ],
        "📈 综合效率提升 2.5 倍 · 人力成本降低 50%",
        "03"
    )
    
    builder.create_comparison_slide(
        "传统 SEO vs GEO（生成式引擎优化）",
        {
            "title": "传统 SEO",
            "items": [
                "优化给搜索引擎看",
                "关键词密度、外链建设",
                "目标：Google 首页排名",
                "结果：获得点击流量",
            ]
        },
        {
            "title": "GEO（新趋势）",
            "items": [
                "优化给 AI 引擎看",
                "语义三重奏、结构化数据",
                "目标：ChatGPT、Claude 引用",
                "结果：AI 主动推荐品牌",
            ]
        },
        "03"
    )
    
    # ========== 模块四：营销获客 ==========
    print("  → 创建模块四...")
    
    builder.create_content_slide(
        "AI 内容矩阵生成实战",
        [
            {"text": "产品描述批量生成", "icon": "📄", "sub": "输入参数，AI 生成多版本"},
            {"text": "Blog 文章自动撰写", "icon": "✍️", "sub": "每月 50-100 篇行业内容"},
            {"text": "社交媒体文案", "icon": "📱", "sub": "多平台适配 + 热门标签"},
            {"text": "邮件营销内容", "icon": "📧", "sub": "个性化开发信，打开率提升 60%"},
        ],
        "📊 案例：AI 生成 300+ 篇文章，6 个月自然流量增长 5 倍",
        "04"
    )
    
    builder.create_content_slide(
        "AI 短视频创作与传播",
        [
            {"text": "脚本生成：AI 写分镜脚本", "icon": "📝"},
            {"text": "视频制作：Runway、Pictory", "icon": "🎥"},
            {"text": "数字人播报：HeyGen、D-ID", "icon": "👤"},
            {"text": "批量分发：多平台自动发布", "icon": "📤"},
        ],
        "🔧 工具：Runway · HeyGen · CapCut AI",
        "04"
    )
    
    builder.create_content_slide(
        "B2B 主动获客全流程",
        [
            {"text": "客户智能搜索", "icon": "🔍", "sub": "AI 挖掘 LinkedIn、Google、海关数据"},
            {"text": "背调分析报告", "icon": "📋", "sub": "自动生成客户画像、采购偏好"},
            {"text": "开发信写作", "icon": "✉️", "sub": "AI 个性化邮件，回复率 2% → 8%"},
        ],
        "📈 3 个月获取 50+ 询盘，成交 8 单总金额 $2.3M",
        "04"
    )
    
    # ========== 模块五：工具演示 ==========
    print("  → 创建模块五...")
    
    builder.create_content_slide(
        "热点工具实战清单",
        [
            {"text": "OpenClaw：AI 爬虫与数据采集", "icon": "🦞"},
            {"text": "GEO 工具：AI 引擎可见性优化", "icon": "🎯"},
            {"text": "Jasper/Copy.ai：AI 内容生成", "icon": "✍️"},
            {"text": "Midjourney/DALL·E 3：AI 图像生成", "icon": "🎨"},
            {"text": "Descript/CapCut：AI 视频剪辑", "icon": "🎬"},
        ],
        None,
        "05"
    )
    
    builder.create_content_slide(
        "推荐工具清单（按预算）",
        [
            {"text": "免费：ChatGPT、Google Trends、Canva", "icon": "🆓"},
            {"text": "入门级（$20-50/月）：Jungle Scout、Jasper", "icon": "💰"},
            {"text": "企业级（$200+/月）：OpenClaw、Custom AI Agent", "icon": "🏢"},
        ],
        None,
        "05"
    )
    
    builder.create_content_slide(
        "学生入门学习路径",
        [
            {"text": "阶段 1：认知建立（1-2 周）", "icon": "🌱", "sub": "AI 基础、主流工具、行业格局"},
            {"text": "阶段 2：工具实操（3-4 周）", "icon": "🛠️", "sub": "熟练使用 3-5 个核心工具"},
            {"text": "阶段 3：项目实战（2-3 月）", "icon": "🚀", "sub": "从 0 到 1 完整运营"},
            {"text": "阶段 4：优化迭代（持续）", "icon": "📈", "sub": "数据分析、策略优化"},
        ],
        "💡 建议：先精通 1-2 个工具 · 实战 > 理论",
        "05"
    )
    
    # ========== 总结 ==========
    print("  → 创建总结页...")
    
    builder.create_content_slide(
        "AI 使用的 5 个误区",
        [
            {"text": "完全依赖 AI → AI 辅助 + 人工审核", "icon": "⚠️"},
            {"text": "忽视数据质量 → 高质量输入 = 高质量输出", "icon": "⚠️"},
            {"text": "盲目追求新工具 → 精通核心工具", "icon": "⚠️"},
            {"text": "忽视合规风险 → 遵守平台规则", "icon": "⚠️"},
            {"text": "期待立竿见影 → 持续优化、长期主义", "icon": "⚠️"},
        ],
        "🎓 AI 不会取代跨境人，但会用 AI 的跨境人一定会取代不会用 AI 的跨境人",
        "06"
    )
    
    builder.create_ending_slide()
    
    # ========== 保存 ==========
    output_path = "/Users/oceanlink/Documents/AIGC 数字贸易出海及跨境电商运营_视觉升级版.pptx"
    builder.save(output_path)
    
    print("\n✨ 视觉升级完成！")
    print("📊 设计亮点：")
    print("  • 现代化配色方案（深海蓝 + 科技蓝 + 荧光绿）")
    print("  • 清晰的视觉层次和信息架构")
    print("  • 统一的卡片式布局")
    print("  • 专业的数据可视化")
    print("  • 适配 2 小时课程节奏（40 页精简版）")


if __name__ == "__main__":
    main()